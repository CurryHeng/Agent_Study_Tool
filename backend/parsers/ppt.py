"""PPT（.pptx）解析器：每页幻灯片一个章节，正文中的教材式小节编号提升为子章节。"""
from pptx import Presentation

from parsers.base import ParsedDocument, Section
from parsers.text_utils import _clean_kp_name, slide_subheading_level


class PptParser:
    source_type = "ppt"

    def parse(self, file_path: str) -> ParsedDocument:
        prs = Presentation(file_path)
        sections: list[Section] = []

        for index, slide in enumerate(prs.slides, start=1):
            title_shape = slide.shapes.title
            title = title_shape.text.strip() if title_shape else ""
            title_el = title_shape._element if title_shape is not None else None
            root = Section(title=title or f"第 {index} 页", level=1, paragraphs=[])
            sections.append(root)
            current = root
            for shape in slide.shapes:
                if shape._element is title_el or not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in p.runs).strip()
                    if not text:
                        continue
                    # 正文行若为教材式小节标题（第X章 / x.y 编号），提升为幻灯片子章节
                    if slide_subheading_level(text) is not None:
                        clean = _clean_kp_name(text)
                        if clean is not None:
                            current = Section(title=clean, level=2, paragraphs=[])
                            sections.append(current)
                            continue
                    current.paragraphs.append(text)

        return ParsedDocument(title="", source_type=self.source_type, sections=sections)
