"""文档解析器测试（markdown / word / ppt + 类型识别）。"""
from parsers.factory import detect_type, parse_file


def test_detect_type():
    assert detect_type("a.md") == "markdown"
    assert detect_type("a.markdown") == "markdown"
    assert detect_type("a.txt") == "markdown"
    assert detect_type("a.pdf") == "pdf"
    assert detect_type("a.docx") == "word"
    assert detect_type("a.pptx") == "ppt"
    assert detect_type("a.png") == "image"
    assert detect_type("a.xyz") is None


def test_markdown_parser(tmp_path):
    md = "# 第一章\n\n内容1\n\n## 知识点A\n\n内容2\n"
    path = tmp_path / "outline.md"
    path.write_text(md, encoding="utf-8")

    parsed = parse_file(str(path), "markdown")
    assert parsed.title == "outline"
    assert len(parsed.sections) == 2
    assert parsed.sections[0].title == "第一章"
    assert parsed.sections[0].level == 1
    assert "内容1" in parsed.sections[0].paragraphs
    assert parsed.sections[1].title == "知识点A"
    assert parsed.sections[1].level == 2


def test_word_parser(tmp_path):
    from docx import Document as Docx

    doc = Docx()
    doc.add_heading("第一章", level=1)
    doc.add_paragraph("段落1")
    doc.add_heading("知识点A", level=2)
    doc.add_paragraph("段落2")
    path = tmp_path / "outline.docx"
    doc.save(str(path))

    parsed = parse_file(str(path), "word")
    titles = [s.title for s in parsed.sections]
    assert "第一章" in titles
    assert "知识点A" in titles


def test_ppt_parser(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "知识点A"
    path = tmp_path / "outline.pptx"
    prs.save(str(path))

    parsed = parse_file(str(path), "ppt")
    assert len(parsed.sections) == 1
    assert parsed.sections[0].title == "知识点A"


def test_ppt_parser_real_courseware(tmp_path):
    """真实课件实测（#63）：标题不泄漏进正文，小节编号提升为子章节。"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "第1章 概述"
    body = slide.placeholders[1].text_frame
    first = True
    for line in ("1.1 定义", "定义的内容。", "1.2 历史", "历史的内容。"):
        para = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        para.text = line
        para.font.size = Pt(14)
    path = tmp_path / "course.pptx"
    prs.save(str(path))

    parsed = parse_file(str(path), "ppt")
    by_title = {s.title: s for s in parsed.sections}
    # 标题占位符文本只出现在章节名，不泄漏进任何段落
    assert "第1章 概述" in by_title
    all_paras = [p for s in parsed.sections for p in s.paragraphs]
    assert "第1章 概述" not in all_paras
    # 教材式小节编号提升为 level 2 子章节，正文归入对应小节
    assert by_title["1.1 定义"].level == 2
    assert "定义的内容。" in by_title["1.1 定义"].paragraphs
    assert by_title["1.2 历史"].level == 2
    assert "历史的内容。" in by_title["1.2 历史"].paragraphs


def test_text_utils_split():
    from parsers.text_utils import split_text_to_sections

    text = "第一章 函数与极限\n函数的定义\n1.1 极限\n极限的定义\n第二章 导数\n"
    parsed = split_text_to_sections(text, "pdf")
    titles = [s.title for s in parsed.sections]
    assert "第一章 函数与极限" in titles
    assert "1.1 极限" in titles
    assert "第二章 导数" in titles


def test_text_utils_year_leading_sentence_is_not_heading():
    """真实 PDF 实测（#63）：以年份/数字开头的正文句不被误判为章节标题。"""
    from parsers.text_utils import split_text_to_sections

    text = (
        "第1章 绪论\n"
        "1956 年达特茅斯会议标志着人工智能学科正式诞生。\n"
        "1.1 研究背景\n"
        "2 到 3 年内该领域快速发展。\n"
        "1 Introduction\n"
        "普通正文。\n"
    )
    parsed = split_text_to_sections(text, "pdf")
    titles = [s.title for s in parsed.sections]
    assert titles == ["第1章 绪论", "1.1 研究背景"]
    # 年份句归入所在章节的正文，而不是变成标题
    assert "1956 年达特茅斯会议标志着人工智能学科正式诞生。" in parsed.sections[0].paragraphs
    assert "2 到 3 年内该领域快速发展。" in parsed.sections[1].paragraphs
